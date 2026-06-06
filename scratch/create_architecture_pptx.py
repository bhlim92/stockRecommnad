import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    BG_COLOR = RGBColor(13, 17, 38)        # #0D1126
    CARD_BG = RGBColor(21, 28, 60)         # #151C3C
    CYAN = RGBColor(0, 242, 254)           # #00F2FE
    PURPLE = RGBColor(185, 0, 255)         # #B900FF
    GREEN = RGBColor(46, 213, 115)         # #2ED573
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(160, 174, 192)

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_title(slide, text, subtitle_text=None):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Malgun Gothic"
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRAY
            p2.space_before = Pt(5)

    def draw_node(slide, left, top, width, height, title, desc, border_color=CYAN):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = "Malgun Gothic"
            p2.font.size = Pt(10)
            p2.font.color.rgb = GRAY
            p2.space_before = Pt(4)
            p2.alignment = PP_ALIGN.CENTER
        return shape

    def draw_arrow(slide, left, top, width, height, color=GRAY, pointing_right=True):
        shape_type = MSO_SHAPE.RIGHT_ARROW if pointing_right else MSO_SHAPE.DOWN_ARROW
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ==========================================
    # Slide 1: Cover Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    # Decorative background circles/shapes
    decor1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
    decor1.fill.solid()
    decor1.fill.fore_color.rgb = RGBColor(18, 30, 70)
    decor1.line.fill.background()
    
    decor2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
    decor2.fill.solid()
    decor2.fill.fore_color.rgb = RGBColor(25, 20, 65)
    decor2.line.fill.background()

    # Title Card
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    
    p_main = tf.paragraphs[0]
    p_main.text = "Stock Discovery & Portfolio Rebalancing System"
    p_main.font.name = "Calibri"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = CYAN
    
    p_sub = tf.add_paragraph()
    p_sub.text = "전체 시스템 아키텍처 및 연동 구조도"
    p_sub.font.name = "Malgun Gothic"
    p_sub.font.size = Pt(28)
    p_sub.font.bold = True
    p_sub.font.color.rgb = WHITE
    p_sub.space_before = Pt(10)

    p_info = tf.add_paragraph()
    p_info.text = "Backend (FastAPI, Threading) & Frontend (Vanilla Web, Polling) Integration Plan"
    p_info.font.name = "Calibri"
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = GRAY
    p_info.space_before = Pt(20)

    # ==========================================
    # Slide 2: Daily Batch Pipeline Architecture
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_title(slide2, "1. 일별 자동화 파이프라인 구조 (Daily Batch Pipeline)", "스케줄러에 의한 데이터 수집, 분석 및 보고서 발행 흐름")

    # Column 1: Trigger & Executor
    draw_node(slide2, Inches(0.8), Inches(1.8), Inches(2.2), Inches(1.2), "작업 스케줄러\n(Windows Task Scheduler)", "매일 오전 지정 시간 실행", CYAN)
    draw_arrow(slide2, Inches(1.7), Inches(3.2), Inches(0.4), Inches(0.6), CYAN, pointing_right=False)
    draw_node(slide2, Inches(0.8), Inches(4.0), Inches(2.2), Inches(1.2), "실행 스크립트\n(run_pipeline.bat)", "가상환경(venv) 활성화 및 라이브러리 검증", CYAN)

    # Arrow to Main Orchestrator
    draw_arrow(slide2, Inches(3.2), Inches(4.3), Inches(0.7), Inches(0.3), CYAN)

    # Column 2: Orchestrator
    draw_node(slide2, Inches(4.1), Inches(3.5), Inches(2.4), Inches(2.0), "메인 오케스트레이터\n(main.py)", "1. 환경 변수 적재\n2. 데이터 수집 기동\n3. Gemini 분석 요청\n4. 자산 비중 계산\n5. 결과 업로드 제어", PURPLE)

    # Data Source boxes connected to Main (top-down or left-right)
    draw_node(slide2, Inches(7.0), Inches(1.8), Inches(2.3), Inches(1.1), "시세 데이터 수집\n(yfinance / FDR)", "Watchlist 주식 및 거시 채권/환율 시세 버퍼링", GREEN)
    draw_node(slide2, Inches(7.0), Inches(3.2), Inches(2.3), Inches(1.1), "뉴스 RSS 수집\n(NewsFetcher)", "S&P500/KOSPI 거시 뉴스 헤드라인 검색", GREEN)
    draw_node(slide2, Inches(7.0), Inches(4.6), Inches(2.3), Inches(1.1), "유튜브 자막 요약\n(YouTubeSummarizer)", "전문가 채널의 영상 스크랩 및 텍스트 추출", GREEN)

    # Arrows from Orchestrator to Data Sources
    draw_arrow(slide2, Inches(6.6), Inches(2.2), Inches(0.3), Inches(0.2), GRAY)
    draw_arrow(slide2, Inches(6.6), Inches(3.6), Inches(0.3), Inches(0.2), GRAY)
    draw_arrow(slide2, Inches(6.6), Inches(5.0), Inches(0.3), Inches(0.2), GRAY)

    # Output boxes connected from Data Sources / Orchestrator
    draw_node(slide2, Inches(10.0), Inches(2.2), Inches(2.5), Inches(1.2), "Gemini AI 모델\n(Recommender)", "수집된 텍스트 컨텍스트 요약 및 최종 리포트 마크다운 빌드", PURPLE)
    draw_node(slide2, Inches(10.0), Inches(4.2), Inches(2.5), Inches(1.2), "구글 드라이브 업로드\n(Google Drive API)", "최종 Investment_Report.md 파일을 Drive 전용 폴더에 보관", GREEN)

    draw_arrow(slide2, Inches(9.4), Inches(2.7), Inches(0.5), Inches(0.2), GRAY)
    draw_arrow(slide2, Inches(9.4), Inches(4.7), Inches(0.5), Inches(0.2), GRAY)

    # ==========================================
    # Slide 3: Web Server & Screener Polling Architecture
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_title(slide3, "2. 웹 서버 및 스크리너 실시간 연동 (Web Server & Screener)", "클라이언트 브라우저와 백엔드 API 간 실시간 폴링 및 리더보드 최적화 구조")

    # Front-end Section (Left Card Frame)
    fe_frame = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8))
    fe_frame.fill.solid()
    fe_frame.fill.fore_color.rgb = CARD_BG
    fe_frame.line.color.rgb = CYAN
    fe_frame.line.width = Pt(1.5)
    
    # Label for Front-end
    fe_lbl = slide3.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(4.5), Inches(0.5))
    fe_lbl.text_frame.text = "FRONTEND (웹 브라우저 클라이언트)"
    fe_lbl.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    fe_lbl.text_frame.paragraphs[0].font.size = Pt(14)
    fe_lbl.text_frame.paragraphs[0].font.bold = True
    fe_lbl.text_frame.paragraphs[0].font.color.rgb = CYAN
    fe_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    draw_node(slide3, Inches(1.2), Inches(2.5), Inches(3.7), Inches(1.1), "vanilla UI (screener.html)", "사용자 설정 모달, 모델 제어 및 탭 네비게이션", WHITE)
    draw_node(slide3, Inches(1.2), Inches(3.8), Inches(3.7), Inches(1.1), "실시간 리더보드 (Top 15)", "스캔 진행 중 분석 완료된 상위 15개 종목 점수 기준 실시간 정렬 렌더링", CYAN)
    draw_node(slide3, Inches(1.2), Inches(5.1), Inches(3.7), Inches(1.1), "2.0초 폴링 타이머 (fetchStatus)", "서버 과부하 방지 및 대역폭 최소화를 위한 2초 주기 상태 질의", PURPLE)

    # Bidirectional arrows representing network requests
    draw_arrow(slide3, Inches(5.5), Inches(3.5), Inches(0.8), Inches(0.3), CYAN, pointing_right=True)
    
    arrow_back = slide3.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(5.5), Inches(4.5), Inches(0.8), Inches(0.3))
    arrow_back.fill.solid()
    arrow_back.fill.fore_color.rgb = PURPLE
    arrow_back.line.fill.background()

    # Network Label
    net_lbl = slide3.shapes.add_textbox(Inches(5.3), Inches(3.9), Inches(1.2), Inches(0.5))
    net_lbl.text_frame.text = "HTTPS API / JSON\nSameSite='None' Cookie"
    net_lbl.text_frame.paragraphs[0].font.name = "Calibri"
    net_lbl.text_frame.paragraphs[0].font.size = Pt(9)
    net_lbl.text_frame.paragraphs[0].font.color.rgb = GRAY
    net_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Back-end Section (Right Card Frame)
    be_frame = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.8))
    be_frame.fill.solid()
    be_frame.fill.fore_color.rgb = CARD_BG
    be_frame.line.color.rgb = PURPLE
    be_frame.line.width = Pt(1.5)

    be_lbl = slide3.shapes.add_textbox(Inches(6.5), Inches(1.9), Inches(6.0), Inches(0.5))
    be_lbl.text_frame.text = "BACKEND (FastAPI - Cafe24 VPS 상주형 서버)"
    be_lbl.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    be_lbl.text_frame.paragraphs[0].font.size = Pt(14)
    be_lbl.text_frame.paragraphs[0].font.bold = True
    be_lbl.text_frame.paragraphs[0].font.color.rgb = PURPLE
    be_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    draw_node(slide3, Inches(6.9), Inches(2.5), Inches(2.4), Inches(1.1), "인증 미들웨어", "HTTPS 상에서 SameSite='None' 및 Secure 설정으로 교차 출처 로그인 유지", PURPLE)
    draw_node(slide3, Inches(9.7), Inches(2.5), Inches(2.4), Inches(1.1), "FastAPI 라우터", "중복 API 라우트 제거로 스캔/상태/정지 동작 단일 창구화 완료", PURPLE)
    draw_node(slide3, Inches(6.9), Inches(3.8), Inches(5.2), Inches(1.1), "백그라운드 워커 스레드 (Threading)", "요청 즉시 리턴 후 병렬 스레드 가동 (Vercel에서는 중단되므로 Cafe24 VPS 고정 운영)", GREEN)
    draw_node(slide3, Inches(6.9), Inches(5.1), Inches(2.4), Inches(1.1), "ThreadPoolExecutor (x15)", "15개 스레드로 yfinance 연산 동시 가속화", GREEN)
    draw_node(slide3, Inches(9.7), Inches(5.1), Inches(2.4), Inches(1.1), "데이터베이스 (RDBMS)", "완료 시점에 MariaDB/PostgreSQL에 최종 분석 결과 백업 저장", PURPLE)

    # Save presentation
    output_dir = r"c:\Users\samsung\proj\stockRecommnad\design"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "system_architecture.pptx")
    prs.save(output_path)
    print(f"[SUCCESS] PowerPoint slide deck created successfully: {output_path}")

if __name__ == "__main__":
    create_deck()
